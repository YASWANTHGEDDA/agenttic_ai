# FusedChatbot/server/ai_core_service/modules/content_creation/md_to_office.py
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches # Use a different alias for docx Inches
import logging # Added for better logging

logger = logging.getLogger(__name__)

def parse_slide_chunk(chunk):
    """Parses a single slide's text block to extract its components."""
    slide_title = "Untitled Slide" # Default if no specific title format is found
    
    # Try to match "### Slide X: Title" or "### Title"
    # Make title regex more flexible: allow any characters after "### " until newline
    title_match = re.match(r"###\s*(?:Slide \d+[:\s]*)?(.*)", chunk)
    if title_match:
        slide_title = title_match.group(1).strip()
    elif chunk.strip().startswith("###"): # Fallback if only "### Some Text"
        slide_title = chunk.strip()[4:].strip()


    # Regex to capture content between **Slide Text Content:** and the next known section or end of chunk
    text_content_match = re.search(
        r"\*\*Slide Text Content:\*\*(.*?)(?=\n\s*\*\*Image Prompt:\*\*|\n\s*\*\*Author Notes for Slide|\Z)", 
        chunk, 
        re.DOTALL | re.IGNORECASE # Added IGNORECASE for flexibility
    )
    raw_text_content = text_content_match.group(1).strip() if text_content_match else ""

    image_prompt_match = re.search(
        r"\*\*Image Prompt:\*\*(.*?)(?=\n\s*\*\*Author Notes for Slide|\Z)", # Removed one lookahead alternative
        chunk,
        re.DOTALL | re.IGNORECASE
    )
    raw_image_prompt = image_prompt_match.group(1).strip() if image_prompt_match else ""
    
    author_notes_match = re.search(
        r"\*\*Author Notes for Slide(?: \d+)?:*\*\*(.*)", # Made slide number optional and colon optional
        chunk, 
        re.DOTALL | re.IGNORECASE
    )
    raw_author_notes = author_notes_match.group(1).strip() if author_notes_match else ""
    
    # If no specific sections found, but there's a title, assume rest of content is 'text_content'
    if not raw_text_content and not raw_image_prompt and not raw_author_notes and slide_title != "Untitled Slide":
        # Remove title part from chunk to get remaining content
        remaining_content = chunk
        if title_match:
            remaining_content = chunk[title_match.end():].strip()
        elif chunk.strip().startswith("###"):
            remaining_content = chunk[chunk.find('\n')+1:].strip() if '\n' in chunk else ''
        
        if remaining_content:
            raw_text_content = remaining_content

    return {
        "title": slide_title,
        "text_content": raw_text_content,
        "image_prompt": raw_image_prompt,
        "author_notes": raw_author_notes
    }

def refined_parse_markdown(markdown_content):
    """Splits the markdown content into slide chunks and parses each chunk."""
    slides_data = []
    
    # Handle empty or whitespace-only input
    if not markdown_content or not markdown_content.strip():
        logger.warning("Markdown content is empty or whitespace only. No slides to parse.")
        return []

    # Regex to split by '---' on its own line, possibly surrounded by whitespace lines.
    # It looks for '---' that is either at the start of a new potential slide definition
    # (i.e., followed by '###') or just as a separator.
    # Using a simpler split and then filtering empty chunks might be more robust.
    # Split by '---' possibly surrounded by newlines.
    raw_chunks = re.split(r'\n\s*---\s*\n', markdown_content)
    
    valid_slide_count = 0
    for chunk_idx, chunk in enumerate(raw_chunks):
        stripped_chunk = chunk.strip()
        if not stripped_chunk: # Skip empty chunks resulting from multiple '---' or leading/trailing '---'
            continue
        
        # A chunk is considered a potential slide if it starts with "### " (our title convention)
        # or if it's the first chunk and contains content.
        if stripped_chunk.startswith("### ") or (valid_slide_count == 0 and stripped_chunk):
            data = parse_slide_chunk(stripped_chunk)
            # Only add if there's a meaningful title or some content
            if data["title"] != "Untitled Slide" or data["text_content"] or data["image_prompt"] or data["author_notes"]:
                slides_data.append(data)
                valid_slide_count += 1
            else:
                logger.debug(f"Skipping chunk {chunk_idx+1} as it parsed to no meaningful content: '{stripped_chunk[:50]}...'")
        else:
            logger.debug(f"Skipping chunk {chunk_idx+1} as it doesn't start with '### ' and isn't the first potential content: '{stripped_chunk[:50]}...'")
            
    if not slides_data and markdown_content.strip(): # If no '---' and no '###', treat whole content as one slide
        logger.info("No '---' or '### ' slide markers found, treating entire content as a single slide.")
        data = parse_slide_chunk(markdown_content.strip()) # Parse the whole content
        if data["title"] != "Untitled Slide" or data["text_content"] or data["image_prompt"] or data["author_notes"]:
            slides_data.append(data)

    logger.info(f"Parsed {len(slides_data)} slides from Markdown.")
    return slides_data


def add_text_to_shape_with_markdown(text_frame, markdown_text, is_title=False, default_font_size_pt=16):
    """Adds markdown-formatted text to a PowerPoint text frame."""
    text_frame.clear() 
    text_frame.word_wrap = True # Enable word wrap

    # Define font sizes
    title_font_size = Pt(32) if is_title else default_font_size_pt 
    content_font_size = default_font_size_pt 
    # Define colors (black background means white text)
    text_color = RGBColor(255, 255, 255) # White

    lines = markdown_text.strip().split('\n')
    if not lines or (len(lines) == 1 and not lines[0].strip()): # Handle empty or whitespace only
        p = text_frame.add_paragraph()
        p.text = " " # Add a non-breaking space or just a space to keep placeholder
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = title_font_size if is_title else content_font_size
        run.font.color.rgb = text_color
        return

    for line_idx, line_text in enumerate(lines):
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT # Default alignment

        # Handle bullet points and indentation
        # Regex: ^(\s*)(?:([\*\-+]|\d+\.)\s+)?(.*)
        # Group 1: Leading spaces (for indentation)
        # Group 2: Bullet character (*, -, +) or number for ordered list (e.g., 1.)
        # Group 3: Actual content of the line
        bullet_match = re.match(r'^(\s*)(?:([\*\-+]|\d+\.)\s+)?(.*)', line_text)
        
        leading_spaces = ""
        bullet_char = None
        content_after_bullet = line_text.strip() # Default to stripped line if no bullet

        if bullet_match:
            leading_spaces = bullet_match.group(1) or ""
            bullet_char = bullet_match.group(2) # This will be None if no bullet/number
            content_after_bullet = bullet_match.group(3).strip()

        # Determine indent level (e.g., 2 spaces per level)
        # python-pptx uses levels 0-8 for bullets
        indent_level = min(len(leading_spaces) // 2, 8) 
        if bullet_char: # If it's a bulleted or numbered list item
            p.level = indent_level
        # else: p.level remains 0 (normal paragraph)

        # Handle bold (**) and italic (*) using simple splitting for now
        # More robust would be a proper Markdown parsing library segment by segment
        # For simplicity, we'll split by bold/italic markers
        # Regex: (\*{1,2})([^*]+)\1 -> *italic* or **bold**
        # This is complex to do perfectly with simple regex splits for nested/mixed styles.
        # We'll handle simple bold for now.
        
        segments = re.split(r'(\*\*.*?\*\*)', content_after_bullet) # Split by bold

        if not any(s.strip() for s in segments): # If line is empty after bullet processing
            if bullet_char: # If it was just a bullet, add a placeholder space
                 run = p.add_run(); run.text = " "
                 run.font.size = content_font_size; run.font.color.rgb = text_color
            # If not a bullet and empty, the empty paragraph p is enough for a blank line
            continue


        for segment in segments:
            if not segment: continue
            run = p.add_run()
            if segment.startswith("**") and segment.endswith("**") and len(segment) > 4:
                run.text = segment[2:-2]
                run.font.bold = True
            else:
                run.text = segment
            
            run.font.size = title_font_size if is_title else content_font_size
            run.font.color.rgb = text_color
            run.font.name = 'Calibri' # Or your preferred font

# *** RENAMED FUNCTION ***
def create_ppt_from_markdown(markdown_content: str, output_dir: str, filename: str) -> str:
    """
    Creates a PowerPoint presentation from parsed slide data.
    This is the function called by app.py.

    Args:
        markdown_content (str): The raw markdown input.
        output_dir (str): Directory to save the PPTX file.
        filename (str): Name of the PPTX file.

    Returns:
        str: Full path to the created PPTX file.
    """
    logger.info(f"Creating PPT: Output Dir='{output_dir}', Filename='{filename}'")
    os.makedirs(output_dir, exist_ok=True)
    output_filepath = os.path.join(output_dir, filename)

    prs = Presentation()
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(16) 
    prs.slide_height = Inches(9)

    slides_data = refined_parse_markdown(markdown_content)

    if not slides_data:
        logger.warning(f"No slides parsed from markdown. Creating a PPT with a 'No Content' slide: {filename}")
        # Add a single slide indicating no content was parsed
        slide_layout = prs.slide_layouts[5] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background; fill = background.fill
        fill.solid(); fill.fore_color.rgb = RGBColor(0, 0, 0) # Black background
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = "No slide content found or parsed from the provided Markdown."
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        prs.save(output_filepath)
        return output_filepath

    for slide_idx, slide_info in enumerate(slides_data):
        slide_layout = prs.slide_layouts[5] # Using a blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set slide background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0) # Black

        # Add Title
        # Centered title, adjust dimensions as needed
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = prs.slide_width - Inches(1.0)
        title_height = Inches(1.2) # Increased height for potentially longer titles
        
        title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_shape.text_frame
        add_text_to_shape_with_markdown(title_tf, slide_info.get("title", f"Slide {slide_idx+1}"), is_title=True, default_font_size_pt=Pt(28))
        # Center title text (optional)
        for para in title_tf.paragraphs: # Ensure all paragraphs in title are centered
            para.alignment = PP_ALIGN.CENTER


        # Add Text Content (if any)
        text_content = slide_info.get("text_content", "")
        if text_content.strip():
            content_left = Inches(0.75)
            content_top = Inches(1.8) # Below title
            content_width = prs.slide_width - Inches(1.5) # Full width content area
            content_height = prs.slide_height - Inches(2.5) # Remaining height
            
            content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            add_text_to_shape_with_markdown(content_shape.text_frame, text_content, default_font_size_pt=Pt(18))

        # Add Author Notes (if any) to the notes slide part
        author_notes = slide_info.get("author_notes", "")
        if author_notes.strip():
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = author_notes
            
        # Placeholder for image_prompt handling - this is complex
        # image_prompt = slide_info.get("image_prompt", "")
        # if image_prompt.strip():
        #     # 1. Call an image generation API with image_prompt
        #     # 2. Download the generated image to a temp path
        #     # 3. Add the image to the slide:
        #     #    img_path = "path/to/generated_image.png"
        #     #    left_img, top_img, width_img, height_img = Inches(8.5), Inches(1.5), Inches(7), Inches(6.5) # Example position
        #     #    slide.shapes.add_picture(img_path, left_img, top_img, width=width_img, height=height_img)
        #     #    os.remove(img_path) # Clean up temp image
        #     pass


    prs.save(output_filepath)
    logger.info(f"PPTX file '{output_filepath}' created successfully with {len(slides_data)} slides.")
    return output_filepath


# --- DOCX Creation Functions ---
def add_markdown_line_to_docx(doc_parent, markdown_line):
    """Adds a single line of markdown-processed text to a docx document/parent (doc or table cell)."""
    # Simplified: handles bold and basic bullets. More complex MD needs a real parser.
    stripped_line = markdown_line.strip() # Process stripped line for content
    original_leading_spaces = len(markdown_line) - len(markdown_line.lstrip(' '))

    if not stripped_line and markdown_line == "": # An actual empty line
        doc_parent.add_paragraph("") # Preserve blank line
        return
    elif not stripped_line: # Line was only whitespace
        doc_parent.add_paragraph("")
        return

    p = doc_parent.add_paragraph()
    
    # Basic bullet handling
    # Regex to match optional leading spaces, then a bullet/number, then content
    bullet_match = re.match(r'^(\s*)([\*\-+]|\d+\.)\s+(.*)', markdown_line) # Match on original line

    if bullet_match:
        spaces_before_bullet = bullet_match.group(1)
        # bullet_char = bullet_match.group(2) # Not used directly, but indicates list style
        content_after_bullet = bullet_match.group(3)
        
        # python-docx list styles are usually 'ListBullet' or 'ListNumber'
        # Indentation needs to be handled by paragraph format if not using multi-level list styles
        p.style = 'ListBullet' # Apply basic bullet style
        
        # Indentation based on leading spaces before the bullet
        indent_level = len(spaces_before_bullet) // 2 # Assuming 2 spaces per indent
        if indent_level > 0:
            p.paragraph_format.left_indent = DocxInches(0.25 * indent_level)
        
        text_to_format = content_after_bullet # Format content after bullet
    else:
        text_to_format = stripped_line # Format the stripped line if not a bullet


    # Basic bold handling: split by **
    segments = re.split(r'(\*\*.*?\*\*)', text_to_format)
    for segment in segments:
        if not segment: continue # Skip empty segments from split
        run = p.add_run()
        if segment.startswith("**") and segment.endswith("**") and len(segment) > 4: # Check length to avoid empty bold
            run.text = segment[2:-2]
            run.font.bold = True
        else:
            run.text = segment

# Renamed function for clarity and consistency (though app.py might call the old name)
# For now, let's keep the name app.py expects for DOCX
def create_doc_from_markdown(markdown_content: str, output_dir: str, filename: str, content_key: str) -> str:
    """
    Creates a Word document from parsed slide data for a specific content key.
    `app.py` calls this function.

    Args:
        markdown_content (str): Raw Markdown input.
        output_dir (str): Directory to save the DOCX file.
        filename (str): Name of the DOCX file.
        content_key (str): The key in slide_data dict to extract content from
                           (e.g., "text_content", "image_prompt", "author_notes").

    Returns:
        str: Full path to the created DOCX file.
    """
    logger.info(f"Creating DOCX: Output Dir='{output_dir}', Filename='{filename}', Content Key='{content_key}'")
    os.makedirs(output_dir, exist_ok=True)
    output_filepath = os.path.join(output_dir, filename)
    
    doc = Document()
    doc_title = os.path.splitext(filename)[0].replace('_', ' ').title() # Nicer title from filename
    doc.add_heading(doc_title, level=0)
    doc.add_paragraph(f"Content extracted for key: '{content_key}'\n")

    slides_data = refined_parse_markdown(markdown_content)
    
    if not slides_data:
        doc.add_paragraph(f"[No slide data parsed from Markdown to extract for '{content_key}']")
        doc.save(output_filepath)
        logger.warning(f"DOCX file '{output_filepath}' created (empty content due to no slide data parsed).")
        return output_filepath

    content_found_for_key = False
    for i, slide_info in enumerate(slides_data):
        slide_title = slide_info.get('title', f"Untitled Slide {i+1}")
        doc.add_heading(f"From Slide: {slide_title}", level=1)
        
        content_to_add = slide_info.get(content_key, "")

        if not content_to_add or not content_to_add.strip():
            doc.add_paragraph(f"[No content found for key '{content_key}' in this slide.]")
        else:
            content_found_for_key = True
            lines = content_to_add.strip().split('\n')
            for line_text in lines:
                add_markdown_line_to_docx(doc, line_text)
        
        if i < len(slides_data) - 1:
            doc.add_paragraph("\n" + "_" * 40 + "\n") # Visual separator

    if not content_found_for_key:
        doc.add_paragraph(f"\n[No content was found for the specified key '{content_key}' across all parsed slides.]")

    doc.save(output_filepath)
    logger.info(f"DOCX file '{output_filepath}' created successfully with content from key '{content_key}'.")
    return output_filepath